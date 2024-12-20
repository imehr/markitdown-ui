import streamlit as st
import tempfile
import os
import re
import zipfile
from docx import Document
from PyPDF2 import PdfReader
import nbformat
from nbconvert import MarkdownExporter
from bs4 import BeautifulSoup
from pptx import Presentation

def split_camel_case(text):
    """Split camelCase and PascalCase text into words"""
    # Handle special cases first
    text = re.sub(r'([a-z])([A-Z][a-z])', r'\1 \2', text)  # camelCase
    text = re.sub(r'([A-Z][a-z])([A-Z][a-z])', r'\1 \2', text)  # PascalCase
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Additional camelCase
    return text

def fix_pdf_spacing(text):
    """Fix common PDF text extraction issues with spacing"""
    if not text:
        return ""

    # Split text into words
    words = text.split()
    processed_words = []
    
    for word in words:
        # Skip short words and special characters
        if len(word) <= 3 or word in ['and', 'the', 'for', 'with']:
            processed_words.append(word)
            continue
            
        # Split long compound words
        if len(word) > 12:  # Likely a compound word
            # Try to split on common patterns
            parts = []
            current_part = ""
            prev_char_type = None
            
            for char in word:
                char_type = 'upper' if char.isupper() else 'lower' if char.islower() else 'other'
                
                # Start new part on type changes
                if prev_char_type and char_type != prev_char_type:
                    if current_part:
                        parts.append(current_part)
                    current_part = char
                else:
                    current_part += char
                
                prev_char_type = char_type
            
            if current_part:
                parts.append(current_part)
            
            # Clean up parts
            cleaned_parts = []
            for part in parts:
                # Handle common word parts
                part = re.sub(r'(?<=[a-z])(?=[A-Z])', ' ', part)  # camelCase
                part = re.sub(r'(?<=[A-Z])(?=[A-Z][a-z])', ' ', part)  # PascalCase
                cleaned_parts.extend(part.split())
            
            processed_words.extend(cleaned_parts)
        else:
            # Handle normal words
            word = re.sub(r'(?<=[a-z])(?=[A-Z])', ' ', word)  # camelCase
            word = re.sub(r'(?<=[A-Z])(?=[A-Z][a-z])', ' ', word)  # PascalCase
            processed_words.extend(word.split())
    
    # Join words back together
    text = ' '.join(processed_words)
    
    # Fix specific patterns
    patterns = [
        (r'(\d):(\d{2})([ap]m)', r'\1:\2 \3'),  # Fix times
        (r'(\d)mins', r'\1 mins'),  # Fix minutes
        (r'(\d{1,2})([ap]m)', r'\1 \2'),  # Fix am/pm
        (r'(\d+)Classification:', r'\1. Classification:'),  # Fix classifications
        (r'\s+([.,!?;:])', r'\1'),  # Fix punctuation spacing
        (r'([.,!?;:])\s*(?=\w)', r'\1 '),
        (r'\s*-\s*', '-'),  # Fix dashes
        (r'(\d)\s+%', r'\1%'),  # Fix percentages
        (r'\s*:\s*', ': '),  # Fix colons
        (r'\s*\(\s*', ' ('),  # Fix parentheses
        (r'\s*\)\s*', ') '),
        (r'•\s*', '• '),  # Fix bullet points
        (r'(\d+\.)\s*', r'\1 '),  # Fix numbered lists
        (r'\s+', ' ')  # Clean up extra spaces
    ]
    
    for pattern, replacement in patterns:
        text = re.sub(pattern, replacement, text)
    
    # Final cleanup of common compound words
    common_compounds = {
        r'AIStrategy': 'AI Strategy',
        r'AIenhanced': 'AI enhanced',
        r'GenAI': 'Gen AI',
        r'userinteraction': 'user interaction',
        r'designvalidation': 'design validation',
        r'usability': 'usability'
    }
    
    for compound, replacement in common_compounds.items():
        text = text.replace(compound, replacement)
    
    return text.strip()

def extract_text_from_pdf(page):
    """Extract text from PDF page while preserving structure"""
    text = page.extract_text()
    if not text:
        return ""
    
    # Split into lines and process each line
    lines = []
    current_paragraph = []
    
    for line in text.split('\n'):
        # Skip empty lines
        if not line.strip():
            if current_paragraph:
                # Join and clean the current paragraph
                paragraph = ' '.join(current_paragraph)
                cleaned = fix_pdf_spacing(paragraph)
                if cleaned:
                    lines.append(cleaned)
                current_paragraph = []
            continue
        
        # Check if line starts with a bullet point or number
        is_list_item = bool(re.match(r'^\s*(?:[•\-*]|\d+\.|\(\d+\))\s', line.strip()))
        
        # If it's a list item or previous paragraph exists, start a new paragraph
        if is_list_item or (current_paragraph and re.search(r'[.!?]$', current_paragraph[-1])):
            if current_paragraph:
                paragraph = ' '.join(current_paragraph)
                cleaned = fix_pdf_spacing(paragraph)
                if cleaned:
                    lines.append(cleaned)
                current_paragraph = []
        
        # Add line to current paragraph
        current_paragraph.append(line.strip())
    
    # Handle any remaining paragraph
    if current_paragraph:
        paragraph = ' '.join(current_paragraph)
        cleaned = fix_pdf_spacing(paragraph)
        if cleaned:
            lines.append(cleaned)
    
    # Join paragraphs with double newlines
    return '\n\n'.join(lines)

def convert_file_to_markdown(file_path):
    """Convert a single file to markdown based on its extension"""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    try:
        if ext in ['.docx']:
            doc = Document(file_path)
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            return '\n\n'.join(paragraphs)
        
        elif ext == '.pdf':
            reader = PdfReader(file_path)
            text_blocks = []
            for page in reader.pages:
                text = extract_text_from_pdf(page)
                if text:
                    text_blocks.append(text)
            return '\n\n'.join(text_blocks)
        
        elif ext in ['.ppt', '.pptx']:
            prs = Presentation(file_path)
            slides = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append('\n'.join(slide_text))
            return '\n\n'.join(slides)
        
        elif ext == '.ipynb':
            with open(file_path, 'r', encoding='utf-8') as f:
                notebook = nbformat.read(f, as_version=4)
            exporter = MarkdownExporter()
            markdown, _ = exporter.from_notebook_node(notebook)
            return markdown
        
        elif ext in ['.html', '.htm']:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Extract text while preserving some structure
                for br in soup.find_all('br'):
                    br.replace_with('\n')
                return soup.get_text('\n\n')
        
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        else:
            return f"Unsupported file format: {ext}"
    
    except Exception as e:
        return f"Error converting file {os.path.basename(file_path)}: {str(e)}"

def convert_to_markdown(uploaded_file):
    if uploaded_file is None:
        return "Please upload a file"
    
    try:
        # Create a temporary directory for processing
        with tempfile.TemporaryDirectory() as temp_dir:
            # Get file extension
            _, ext = os.path.splitext(uploaded_file.name)
            ext = ext.lower()
            
            # Save uploaded file to temporary directory
            temp_file = os.path.join(temp_dir, uploaded_file.name)
            with open(temp_file, 'wb') as f:
                f.write(uploaded_file.getvalue())
            
            # Handle ZIP files
            if ext == '.zip':
                try:
                    markdown_parts = []
                    with zipfile.ZipFile(temp_file, 'r') as zip_ref:
                        # Extract all files
                        zip_ref.extractall(temp_dir)
                        # Convert each file
                        for root, _, files in os.walk(temp_dir):
                            for file in files:
                                if file != uploaded_file.name:  # Skip the zip file itself
                                    file_path = os.path.join(root, file)
                                    relative_path = os.path.relpath(file_path, temp_dir)
                                    markdown_parts.append(f"# {relative_path}")
                                    markdown_parts.append(convert_file_to_markdown(file_path))
                    
                    return '\n\n'.join(markdown_parts)
                except zipfile.BadZipFile:
                    return "Error: Invalid ZIP file"
            else:
                return convert_file_to_markdown(temp_file)
    
    except Exception as e:
        return f"Error converting file: {str(e)}"

def main():
    st.set_page_config(
        page_title="MarkItDown - File to Markdown Converter",
        page_icon="📝",
        layout="wide"
    )
    
    st.title("MarkItDown - File to Markdown Converter")
    st.write("Convert Documents (PDF, Word, PowerPoint), Code (Jupyter), Web Pages (HTML, TXT), or ZIP archives to Markdown.")
    
    uploaded_files = st.file_uploader(
        "Choose files",
        type=['pdf', 'docx', 'ppt', 'pptx', 'ipynb', 'html', 'htm', 'txt', 'zip'],
        help="Upload files to convert to Markdown. Supported formats: PDF, Word, PowerPoint, Jupyter Notebook, HTML, TXT, and ZIP archives containing these files.",
        accept_multiple_files=True
    )
    
    if uploaded_files:
        st.write(f"Files uploaded: {', '.join(f.name for f in uploaded_files)}")
        
        if st.button("Convert to Markdown", type="primary"):
            with st.spinner('Converting...'):
                markdown_parts = []
                
                # Process each uploaded file
                for uploaded_file in uploaded_files:
                    markdown_parts.append(f"# {uploaded_file.name}")
                    markdown_parts.append(convert_to_markdown(uploaded_file))
                
                markdown_output = '\n\n'.join(markdown_parts)
                
                st.markdown("### Output:")
                st.code(markdown_output, language='markdown')
                
                # Add download button for the converted markdown
                st.download_button(
                    label="Download Markdown",
                    data=markdown_output,
                    file_name="converted_files.md",
                    mime="text/markdown"
                )

if __name__ == "__main__":
    main()
